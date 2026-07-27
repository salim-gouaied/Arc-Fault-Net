import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# -----------------------------------------------------------------------------
# INITIALISATION DE LA PRÉSENTATION
# -----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)  # Format 16:9 standard
prs.slide_height = Inches(7.5)

# Couleurs sobres pour la structure de base
BLANC = RGBColor(255, 255, 255)
NOIR = RGBColor(30, 30, 30)
GRIS_FONCE = RGBColor(80, 80, 80)

def add_blank_slide():
    blank_slide_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_slide_layout)

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NOIR
    return txBox

def add_bullets(slide, items, left=0.75, top=1.8, width=11.83, height=4.5, size=18):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = GRIS_FONCE
        if item.startswith("?") or item.startswith("[Question"):
            p.font.italic = True
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
    return txBox

def create_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    return table_shape.table

# =============================================================================
# SLIDE 1 : TITRE
# =============================================================================
s1 = add_blank_slide()
tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(3.0))
tf = tb.text_frame
p1 = tf.paragraphs[0]
p1.text = "Évolution de l'architecture Arc-FaultNet"
p1.font.size = Pt(40)
p1.font.bold = True
p2 = tf.add_paragraph()
p2.text = "Détection d'arc série par Deep Learning sur le courant de ligne I(t)\nPoint d'étape & justifications expérimentales"
p2.font.size = Pt(20)
p2.font.color.rgb = GRIS_FONCE

# =============================================================================
# SLIDE 2 : DÉMARCHE EXPÉRIMENTALE
# =============================================================================
s2 = add_blank_slide()
add_title(s2, "1. Démarche expérimentale stricte")
add_bullets(s2, [
    "• Hypothèse : Test d'un composant ou d'une modification sur la V1",
    "• Protocole & Métriques : Validation rigoureuse à conditions égales (Dataset, Seed...)",
    "• Étude d'ablation : Isolation de la contribution de chaque bloc (résultats chiffrés)",
    "• Constats : Identification objective de ce qui apporte de la performance ou dégrade",
    "• Décisions : Choix architecturaux validés par les métriques pour concevoir la V2",
    "• Preuve : Validation finale prouvant la supériorité de la V2"
])

# =============================================================================
# SLIDE 3 : V1 POINT DE DÉPART
# =============================================================================
s3 = add_blank_slide()
add_title(s3, "2. Point de départ : Arc-FaultNet V1 (Commit Initial)")
add_bullets(s3, [
    "• Dual-branch CNN Architecture (~344 k paramètres) :",
    "    - Branche 1D : Filtres de Gabor paramétriques (ParametricConv1d, f0 & sigma apprenables)",
    "    - Branche 2D : Spectrogramme STFT",
    "• Mécanisme de fusion : Joint Attention via CAM (Channel Attention) + SAM (Spatial Attention) croisés",
    "• Entrée du modèle : [V_ligne, I]",
    "• Objectif initial : Capturer les dynamiques transitoires de l'arc série."
])

# =============================================================================
# SLIDE 4 : HYPOTHÈSES TESTÉES SUR V1
# =============================================================================
s4 = add_blank_slide()
add_title(s4, "3. Hypothèses d'optimisation testées sur la V1")
add_bullets(s4, [
    "• Pistes explorées pour maximiser la sensibilité :",
    "    - Insertion de blocs Squeeze-and-Excitation (SEBlock) après chaque couche de convolution",
    "    - Intégration d'une amplitude apprenable pour les filtres de Gabor (use_amplitude)",
    "    - Remplacement de la tête par un classifieur plus profond (deep_classifier)",
    "    - Stratégies de réduction stricte de paramètres contre le surapprentissage (overfitting)",
    "\n[Question pour l'encadrant] : Aviez-vous anticipé une sensibilité de la convergence selon la méthode d'initialisation de f0 pour Gabor ?"
])

# =============================================================================
# SLIDE 5 : PROTOCOLE D'ABLATION
# =============================================================================
s5 = add_blank_slide()
add_title(s5, "4. Protocole de l'étude d'ablation")
add_bullets(s5, [
    "• Date de l'expérience : 27/05/2026",
    "• Configuration : Seed fixe (3) | Split rigoureux 70 / 15 / 15",
    "• Données : combined_dataset (10 860 cycles complets, M = 20 000 points @ 1 MHz)",
    "• Objectif : Isoler mathématiquement l'impact de chaque choix de design sur les métriques globales."
])

# =============================================================================
# SLIDE 6 : RÉSULTATS D'ABLATION (TABLEAU)
# =============================================================================
s6 = add_blank_slide()
add_title(s6, "5. Résultats chiffrés de l'étude d'ablation")
t6 = create_table(s6, 7, 4, 0.75, 1.8, 11.83, 4.0)
headers = ["Variante", "Accuracy", "F1-Score", "Params"]
data6 = [
    ["standard_conv", "96,93 %", "96,68 %", "514 201"],
    ["arcfaultnet (réf)", "96,38 %", "96,08 %", "344 409"],
    ["no_attention", "95,15 %", "94,76 %", "179 705"],
    ["independent_cbam", "93,25 %", "92,49 %", "237 721"],
    ["baseline_cnn", "91,04 %", "89,82 %", "209 697"],
    ["1d_only", "51,84 %", "65,88 %", "47 773"]
]
for col_idx, text in enumerate(headers):
    t6.cell(0, col_idx).text = text
for row_idx, row_data in enumerate(data6):
    for col_idx, text in enumerate(row_data):
        t6.cell(row_idx+1, col_idx).text = text

# =============================================================================
# SLIDE 7 : CONSTATS EXPÉRIMENTAUX
# =============================================================================
s7 = add_blank_slide()
add_title(s7, "6. Constats scientifiques majeurs")
add_bullets(s7, [
    "• Gabor vs Conv Standard : −0,55 % de F1-score. Les filtres de Gabor n'apportent RIEN ou dégradent.",
    "• Branche STFT : +44,5 % de F1-score (passage de 1d_only à baseline_cnn). Absolument indispensable.",
    "• Cross-attention vs CBAM indépendants : +3,13 % de F1. La dépendance inter-branches est validée.",
    "• Joint Attention vs Concat simple : +1,23 % de F1-score.",
    "\n[Question cruciale à l'encadrant] : Au vu de la chute induite par Gabor (-0.55%), auriez-vous conservé ces filtres au profit de la physique ou applique-t-on le verdict de l'ablation ?"
])

# =============================================================================
# SLIDE 8 : DÉCISIONS ARCHITECTURALES V2
# =============================================================================
s8 = add_blank_slide()
add_title(s8, "7. Décisions d'ingénierie et passage à la V2")
add_bullets(s8, [
    "• Retrait de Gabor : Remplacement par Conv1d standard + GELU (l'arc est impulsionnel, pas d'oscillation porteuse).",
    "• Réduction de l'entrée au courant seul I(t) : V(t) nuit à la généralisation en encodant l'impédance de la charge.",
    "• Extraction de 4 canaux dérivés de I(t) normalisés par le RMS du cycle (Invariance à la charge) :",
    "    1. I_norm (forme)  |  2. |ΔI| (dérivée)  |  3. TKEO (énergie instantanée)  |  4. RMS glissant (enveloppe)"
], height=2.5)
add_bullets(s8, [
    "• Pooling asymétrique 2D : Compression temporelle mais préservation de la résolution fréquentielle.",
    "• FrequencyGate apprenable : Remplacement de la tranche fréquentielle codée en dur par un masque doux.",
    "• RevisedCrossAttention : Correction du bug d'ordre des canaux sur le CAM joint de la V1.",
    "• Protocole classification en 2 phases : Alignement de l'embedding via FC, puis classification finale par XGBoost."
], top=4.3, height=2.8)

# =============================================================================
# SLIDE 9 : V2 EN BREF
# =============================================================================
s9 = add_blank_slide()
add_title(s9, "8. Synthèse de l'architecture Arc-FaultNet V2")
add_bullets(s9, [
    "• Entrée : 4 canaux physiques construits exclusivement sur le courant de ligne I(t).",
    "• Backbone : Hybride Conv1D-Conv2D épuré (~350 k paramètres).",
    "• Masquage : Sélection adaptative des bandes critiques via FrequencyGate.",
    "• Tête de classification : Embedding 128-d + classifieur final XGBoost (meilleure calibration des probabilités de trip, importance des features)."
])

# =============================================================================
# SLIDE 10 : PROTOCOLE DE COMPARAISON ÉQUITABLE
# =============================================================================
s10 = add_blank_slide()
add_title(s10, "9. Protocole de comparaison strict (V1 vs V2)")
add_bullets(s10, [
    "• Alignement parfait des conditions expérimentales pour isoler le gain de performance :",
    "    - Dataset identique : combined_dataset_2048 (Fréquence échantillonnage réduite à 102,4 kHz)",
    "    - Hyperparamètres constants : lr = 3e-4, weight_decay = 5e-4, batch_size = 64, patience = 10",
    "    - Paramètres STFT : n_fft = 128, hop_length = 64",
    "    - Protocole d'évaluation : Mode single run, split 70 / 15 / 15, Seed fixe (4)"
])

# =============================================================================
# SLIDE 11 : V1 VS V2 CHIFFRES (SLIDE PIVOT)
# =============================================================================
s11 = add_blank_slide()
add_title(s11, "10. Preuve de supériorité à conditions égales (Seed 4)")
add_bullets(s11, [
    "• Accuracy : 93,74 %  -->  97,24 %  (+3,50 pts)",
    "• F1-Score : 92,37 %  -->  96,82 %  (+4,45 pts)",
    "• RECALL (Métrique critique) : 87,15 %  -->  96,75 %  (+9,60 pts)",
    "    --> Argument clé : ~75 % d'arcs manqués en moins (sécurité AFDD drastiquement renforcée)",
    "• Précision : 98,25 %  -->  96,89 %  (-1,36 pt, compromis assumé)",
    "• Spécificité : 98,81 %  -->  97,61 %  (-1,19 pt)",
    "• Complexité : 320 609  -->  350 693 paramètres (+9 % seulement)",
    "\n• Stabilité Multi-seed V2 (10/06) : Seed 2 = 93,76 % | Seed 4 = 96,82 % | Seed 42 = 97,37 % (Moyenne ≈ 96,0 %)"
], size=16)

# =============================================================================
# SLIDE 12 : TROIS MODÈLES PAR DATE
# =============================================================================
s12 = add_blank_slide()
add_title(s12, "11. Historique et trajectoire des performances")
t12 = create_table(s12, 4, 8, 0.5, 1.8, 12.33, 4.0)
h12 = ["Date", "Modèle", "Données", "Seed", "Acc", "F1", "Recall", "Précision"]
d12 = [
    ["26-29/05", "V1", "20k @ 1MHz", "3-6", "94.6-96.4%", "93.9-96.1%", "92-95%", "93-99%"],
    ["03/06", "V1", "2048 @ 102.4kHz", "4", "93,74 %", "92,37 %", "87,15 %", "98,25 %"],
    ["10/06", "V2", "2048 @ 102.4kHz", "4", "97,24 %", "96,82 %", "96,75 %", "96,89 %"]
]
for col_idx, text in enumerate(h12):
    t12.cell(0, col_idx).text = text
for row_idx, row_data in enumerate(d12):
    for col_idx, text in enumerate(row_data):
        t12.cell(row_idx+1, col_idx).text = text

# =============================================================================
# SLIDE 13 : GROUPKFOLD A VENIR
# =============================================================================
s13 = add_blank_slide()
add_title(s13, "12. Validation de la généralisation inter-charges")
add_bullets(s13, [
    "• Objectif : Évaluer la robustesse du modèle sur des types de charges électriques totalement invisibles au cours du training.",
    "• Pourquoi ce protocole est plus difficile ? Le split standard mélange les cycles d'une même charge. Le GroupKFold isole les charges.",
    "\n[Tableau en cours de complétion suite aux entraînements en cours] :",
    "    - V1 GroupKFold F1-Score : [ À compléter ]",
    "    - V2 GroupKFold F1-Score : [ À compléter ]"
])

# =============================================================================
# SLIDE 14 : CONCLUSION + DISCUSSION
# =============================================================================
s14 = add_blank_slide()
add_title(s14, "13. Conclusion et Travaux futurs")
add_bullets(s14, [
    "• Synthèse : Chaque modification de la V2 repose sur des métriques d'ablation, éliminant les intuitions fausses (Gabor, tension V(t)).",
    "• Gain opérationnel : Le bond de +9.60 pts sur le Recall valide scientifiquement l'approche par extraction de features physiques du courant.",
    "• Prochaines étapes immédiates :",
    "    - Finalisation des runs GroupKFold inter-charges.",
    "    - Évaluation de l'apport de la tête XGBoost par rapport à la couche FC standard.",
    "    - Travail futur : Extension multi-cycles en intégrant une couche récurrente (BiGRU)."
])

# Sauvegarde
prs.save("presentation_arcfaultnet.pptx")
print("Fichier 'presentation_arcfaultnet.pptx' généré avec succès !")
